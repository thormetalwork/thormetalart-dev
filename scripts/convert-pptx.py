#!/usr/bin/env python3
"""
convert-pptx.py — Convierte la presentación kickoff PPTX a HTML slides
Ejecutar: python3 scripts/convert-pptx.py
"""

import html
import os
from pptx import Presentation
from pptx.util import Pt

PROJECT_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
INPUT = os.path.join(PROJECT_ROOT, 'docs', 'cliente', 'Portal Cliente', 'thor_kickoff_deck.pptx')
OUTPUT = os.path.join(PROJECT_ROOT, 'portal', 'docs', 'thor_kickoff_deck.html')
CSS_REL = '../css/doc-viewer.css'


def extract_slide_html(slide, idx, total):
    """Extract text from a slide and format as HTML."""
    blocks = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            # Detect heading vs body by font size
            is_heading = False
            is_bold = False
            font_size = None
            for run in para.runs:
                if run.font.size:
                    font_size = run.font.size
                if run.font.bold:
                    is_bold = True

            if font_size and font_size >= Pt(20):
                is_heading = True

            escaped = html.escape(text)
            if is_heading:
                blocks.append(f'<h2>{escaped}</h2>')
            elif is_bold:
                blocks.append(f'<p><strong>{escaped}</strong></p>')
            else:
                blocks.append(f'<p>{escaped}</p>')

    content = '\n      '.join(blocks)
    return f'''    <div class="slide" data-slide="{idx}"{' style="display:block"' if idx == 1 else ''}>
      <div class="slide-num">Slide {idx} / {total}</div>
      {content}
    </div>'''


def main():
    prs = Presentation(INPUT)
    total = len(prs.slides)
    slides_html = '\n'.join(
        extract_slide_html(slide, i + 1, total)
        for i, slide in enumerate(prs.slides)
    )

    full_html = f'''<!DOCTYPE html>
<html lang="es">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>Kickoff Deck — Thor Metal Art</title>
<link rel="preconnect" href="https://fonts.googleapis.com">
<link href="https://fonts.googleapis.com/css2?family=Cormorant+Garamond:ital,wght@0,400;0,600;1,400&family=DM+Sans:opsz,wght@9..40,300;9..40,400;9..40,500&display=swap" rel="stylesheet">
<link rel="stylesheet" href="{CSS_REL}">
<style>
  .slides-container{{position:relative;margin-bottom:32px}}
  .slide{{display:none;background:var(--dark2);border:1px solid var(--border);border-radius:var(--r);padding:40px 36px;min-height:380px}}
  .slide h2{{font-family:'Cormorant Garamond',serif;font-size:28px;color:var(--gold-l);margin:24px 0 16px;line-height:1.2}}
  .slide h2:first-of-type{{margin-top:0}}
  .slide p{{margin:0 0 10px;font-size:15px;line-height:1.7}}
  .slide-num{{font-size:11px;color:var(--dim);letter-spacing:.1em;text-transform:uppercase;margin-bottom:20px}}
  .slide-nav{{display:flex;align-items:center;justify-content:center;gap:16px;margin-top:24px}}
  .slide-btn{{background:var(--dark3);color:var(--text);border:1px solid var(--border);padding:10px 24px;border-radius:8px;cursor:pointer;font:inherit;font-size:14px;transition:all .2s}}
  .slide-btn:hover{{background:rgba(184,134,11,.15);border-color:var(--gold);color:var(--gold)}}
  .slide-btn:disabled{{opacity:.3;cursor:default}}
  .slide-counter{{font-size:14px;color:var(--muted);min-width:60px;text-align:center}}
  .slide-dots{{display:flex;gap:6px;justify-content:center;margin-top:16px}}
  .slide-dot{{width:8px;height:8px;border-radius:50%;background:var(--dark3);border:1px solid var(--border);cursor:pointer;transition:all .2s}}
  .slide-dot.active{{background:var(--gold);border-color:var(--gold)}}
</style>
</head>
<body>

<header>
  <div class="logo">
    <div class="lmark">T</div>
    <div class="ltxt">THOR <span>METAL ART</span></div>
  </div>
  <nav>
    <a href="../">← Volver al Portal</a>
  </nav>
</header>

<main class="doc-container">
  <div class="doc-header">
    <span class="doc-num">Presentación</span>
    <h1>Kickoff Deck — Plan Integral de Presencia Digital</h1>
  </div>

  <div class="slides-container">
{slides_html}
  </div>

  <div class="slide-nav">
    <button class="slide-btn" id="prevBtn">← Anterior</button>
    <span class="slide-counter" id="counter">1 / {total}</span>
    <button class="slide-btn" id="nextBtn">Siguiente →</button>
  </div>
  <div class="slide-dots" id="dots"></div>

  <footer class="doc-footer">
    <a href="../" class="back-link">← Volver al Portal</a>
    <span class="doc-meta">Thor Metal Art — Presentación del Proyecto</span>
  </footer>
</main>

<script>
(function(){{
  const slides = document.querySelectorAll('.slide');
  const total = slides.length;
  let current = 1;
  const counter = document.getElementById('counter');
  const prev = document.getElementById('prevBtn');
  const next = document.getElementById('nextBtn');
  const dotsEl = document.getElementById('dots');

  // Create dots
  for (let i = 1; i <= total; i++) {{
    const dot = document.createElement('span');
    dot.className = 'slide-dot' + (i === 1 ? ' active' : '');
    dot.dataset.slide = i;
    dot.addEventListener('click', function() {{ goTo(parseInt(this.dataset.slide)); }});
    dotsEl.appendChild(dot);
  }}

  function goTo(n) {{
    slides.forEach(s => s.style.display = 'none');
    dotsEl.querySelectorAll('.slide-dot').forEach(d => d.classList.remove('active'));
    current = n;
    slides[current - 1].style.display = 'block';
    dotsEl.querySelector('[data-slide="' + current + '"]').classList.add('active');
    counter.textContent = current + ' / ' + total;
    prev.disabled = current === 1;
    next.disabled = current === total;
  }}

  prev.addEventListener('click', function() {{ if (current > 1) goTo(current - 1); }});
  next.addEventListener('click', function() {{ if (current < total) goTo(current + 1); }});

  // Keyboard navigation
  document.addEventListener('keydown', function(e) {{
    if (e.key === 'ArrowLeft' && current > 1) goTo(current - 1);
    if (e.key === 'ArrowRight' && current < total) goTo(current + 1);
  }});

  goTo(1);
}})();
</script>

</body>
</html>'''

    os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
    with open(OUTPUT, 'w', encoding='utf-8') as f:
        f.write(full_html)

    print(f'✅ PPTX convertido a HTML slides: {OUTPUT}')
    print(f'   Slides: {total}')


if __name__ == '__main__':
    main()
